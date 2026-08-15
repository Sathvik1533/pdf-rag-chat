"""
Universal Document & Data Extractors for Veritas RAG.
Supports PDF, DOCX, PPTX, XLSX, CSV, TSV, JSON, YAML, Source Code, Markdown, and Plain Text.
"""

import io
import csv
import json
import logging
from typing import List, Tuple, Optional
from pathlib import Path

logger = logging.getLogger(__name__)

# Supported file extensions set
SUPPORTED_EXTENSIONS = {
    # Documents
    ".pdf", ".docx", ".doc", ".odt", ".rtf", ".txt", ".md", ".markdown", ".rst", ".log",
    # Presentations
    ".pptx", ".ppt",
    # Spreadsheets & Tabular Data
    ".xlsx", ".xls", ".csv", ".tsv",
    # Structured Data
    ".json", ".jsonl", ".yaml", ".yml", ".xml", ".html", ".htm",
    # Source Code
    ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cpp", ".c", ".h", ".hpp",
    ".go", ".rs", ".sql", ".sh", ".bash", ".zsh", ".css", ".scss", ".php", ".rb", ".swift", ".kt"
}


class ExtractedUnit:
    """Represents a single extracted slice (Page, Slide, Sheet, or Code Block)."""
    def __init__(self, index: int, label: str, text: str):
        self.index = index
        self.label = label
        self.text = text

    def to_tuple(self) -> Tuple[int, str, str]:
        return (self.index, self.label, self.text)


def extract_pdf(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """
    Extract text from PDF page-by-page.
    Tolerates truncated streams, damaged xref tables, compression anomalies,
    and mislabeled files with multi-tier automated recovery.
    """
    import pypdf
    units: List[ExtractedUnit] = []

    # Tier 1: Standard pypdf extraction with lenient parsing
    try:
        reader = pypdf.PdfReader(io.BytesIO(file_bytes), strict=False)
        num_pages = len(reader.pages)
        for idx in range(num_pages):
            try:
                page = reader.pages[idx]
                page_text = page.extract_text() or ""
                cleaned = page_text.strip()
                if cleaned:
                    units.append(ExtractedUnit(index=idx + 1, label=f"Page {idx + 1}", text=cleaned))
            except Exception as pe:
                logger.warning(f"Error on page {idx + 1} of {filename}: {pe}")
                
        if units:
            return units
    except Exception as e:
        logger.warning(f"pypdf reader failed on {filename}: {e}. Activating PDF raw stream recovery...")

    # Tier 2: Raw PDF stream decompresion & operator text extraction
    stream_passages = extract_pdf_raw_streams(file_bytes)
    if stream_passages:
        unit_idx = 1
        chunk_size = 8
        for start in range(0, len(stream_passages), chunk_size):
            batch = stream_passages[start:start + chunk_size]
            units.append(ExtractedUnit(
                index=unit_idx,
                label=f"Page {unit_idx}",
                text="\n\n".join(batch)
            ))
            unit_idx += 1
        if units:
            return units

    # Tier 3: Binary string & UTF-16 stream recovery (in case the file was converted or renamed)
    return extract_binary_doc(file_bytes, filename)


def extract_pdf_raw_streams(file_bytes: bytes) -> List[str]:
    """
    Decompresses flate streams and extracts readable PDF text chunks directly
    even when PDF xref tables or trailers are truncated or damaged.
    """
    import zlib
    import re

    passages = []
    
    # 1. Search for FlateDecode compressed streams
    stream_pattern = re.compile(b"stream[\r\n]+(.*?)[\r\n]+endstream", re.DOTALL)
    for match in stream_pattern.finditer(file_bytes):
        stream_data = match.group(1)
        decompressed = None
        for wbits in [zlib.MAX_WBITS, -zlib.MAX_WBITS, zlib.MAX_WBITS | 32]:
            try:
                decompressed = zlib.decompress(stream_data, wbits)
                break
            except Exception:
                continue
                
        if decompressed:
            # Extract strings between parentheses (text) Tj or [(text)] TJ
            text_matches = re.findall(r'\(([^\)\\]*(?:\\.[^\)\\]*)*)\)\s*(?:Tj|TJ|\')', decompressed.decode("latin-1", errors="ignore"))
            for tm in text_matches:
                cleaned = re.sub(r'\\[nrtbf\\()]', ' ', tm).strip()
                if len(cleaned) > 3 and any(c.isalpha() for c in cleaned):
                    passages.append(cleaned)

    # 2. Extract uncompressed PDF text operators
    try:
        raw_text = file_bytes.decode("latin-1", errors="ignore")
        uncomp_matches = re.findall(r'\(([^\)\\]*(?:\\.[^\)\\]*)*)\)\s*Tj', raw_text)
        for tm in uncomp_matches:
            cleaned = re.sub(r'\\[nrtbf\\()]', ' ', tm).strip()
            if len(cleaned) > 5 and cleaned not in passages:
                passages.append(cleaned)
    except Exception:
        pass

    # Deduplicate and group into coherent paragraphs
    filtered_passages = []
    for p in passages:
        if len(p) > 8 and not any(meta in p for meta in ["Font", "Identity", "ToUnicode"]):
            if p not in filtered_passages:
                filtered_passages.append(p)

    return filtered_passages


def extract_docx(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """
    Extract text from Word .docx document.
    Falls back to binary .doc extraction if it is an older Word 97-2004 binary format.
    """
    try:
        import docx
        doc = docx.Document(io.BytesIO(file_bytes))
        units: List[ExtractedUnit] = []
        
        current_section = []
        section_idx = 1
        word_count = 0
        
        for p in doc.paragraphs:
            txt = p.text.strip()
            if not txt:
                continue
                
            current_section.append(txt)
            word_count += len(txt.split())
            
            # Split roughly every ~350 words or on Heading styles
            is_heading = p.style and p.style.name and "Heading" in p.style.name
            if word_count >= 350 or (is_heading and word_count >= 150):
                units.append(ExtractedUnit(
                    index=section_idx,
                    label=f"Section {section_idx}",
                    text="\n\n".join(current_section)
                ))
                section_idx += 1
                current_section = []
                word_count = 0
                
        if current_section:
            units.append(ExtractedUnit(
                index=section_idx,
                label=f"Section {section_idx}",
                text="\n\n".join(current_section)
            ))
            
        if units:
            return units
            
    except Exception as e:
        logger.info(f"docx XML parser bypassed for {filename}: {e}. Trying binary .doc extractor.")
        
    return extract_binary_doc(file_bytes, filename)


def extract_binary_doc(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """
    Extract readable human text from binary Word 97-2004 (.doc), RTF, or raw binary containers.
    Extracts UTF-16LE (both even and odd alignments) and ASCII text streams while filtering out OLE metadata noise.
    """
    import re
    
    extracted_passages = []

    # 1. Try extracting UTF-16LE strings (checking both even and odd byte alignments)
    for offset in [0, 1]:
        try:
            utf16_text = file_bytes[offset:].decode("utf-16le", errors="ignore")
            utf16_matches = re.findall(r'[A-Za-z0-9][A-Za-z0-9\s,.;:\'"\-()?!/@#%&*+=]{10,}', utf16_text)
            for match in utf16_matches:
                cleaned = " ".join(match.split()).strip()
                # Filter out OLE internal table signatures
                if len(cleaned) > 15 and not any(meta in cleaned for meta in ["WordDocument", "SummaryInformation", "CompObj", "Root Entry"]):
                    if cleaned not in extracted_passages:
                        extracted_passages.append(cleaned)
        except Exception:
            pass

    # 2. Try extracting ASCII / UTF-8 strings
    try:
        ascii_text = file_bytes.decode("utf-8", errors="ignore")
        ascii_matches = re.findall(r'[A-Za-z0-9][A-Za-z0-9\s,.;:\'"\-()?!/@#%&*+=]{10,}', ascii_text)
        for match in ascii_matches:
            cleaned = " ".join(match.split()).strip()
            if len(cleaned) > 15 and not any(meta in cleaned for meta in ["WordDocument", "SummaryInformation", "CompObj", "Root Entry"]):
                if cleaned not in extracted_passages:
                    extracted_passages.append(cleaned)
    except Exception:
        pass

    if not extracted_passages:
        # Fallback to general plain text decoder
        return extract_text_fallback(file_bytes, filename)

    # Group into readable sections of ~350 words
    units: List[ExtractedUnit] = []
    current_section = []
    section_idx = 1
    word_count = 0

    for passage in extracted_passages:
        current_section.append(passage)
        word_count += len(passage.split())
        if word_count >= 350:
            units.append(ExtractedUnit(
                index=section_idx,
                label=f"Section {section_idx}",
                text="\n\n".join(current_section)
            ))
            section_idx += 1
            current_section = []
            word_count = 0

    if current_section:
        units.append(ExtractedUnit(
            index=section_idx,
            label=f"Section {section_idx}",
            text="\n\n".join(current_section)
        ))

    return units


def extract_pptx(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Extract slide text and notes from PowerPoint presentation slide-by-slide."""
    try:
        import pptx
        prs = pptx.Presentation(io.BytesIO(file_bytes))
        units: List[ExtractedUnit] = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            slide_texts = []
            
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            slide_texts.append(txt)
                            
            # Check for speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"[Speaker Notes]: {notes}")
                    
            if slide_texts:
                units.append(ExtractedUnit(
                    index=idx,
                    label=f"Slide {idx}",
                    text="\n".join(slide_texts)
                ))
                
        return units
    except Exception as e:
        logger.warning(f"pptx library error on {filename}: {e}. Falling back to text extraction.")
        return extract_text_fallback(file_bytes, filename)


def extract_xlsx(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Extract Excel sheets formatted as Markdown tables."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
        units: List[ExtractedUnit] = []
        unit_idx = 1
        
        for sheetname in wb.sheetnames:
            ws = wb[sheetname]
            rows = list(ws.iter_rows(values_only=True))
            if not rows:
                continue
                
            # Filter completely empty rows
            non_empty_rows = [
                [str(c) if c is not None else "" for c in row]
                for row in rows
                if any(c is not None and str(c).strip() != "" for c in row)
            ]
            
            if not non_empty_rows:
                continue
                
            headers = non_empty_rows[0]
            data_rows = non_empty_rows[1:]
            
            # Chunk every 30 rows into a separate sheet slice
            chunk_size = 30
            for start in range(0, max(len(data_rows), 1), chunk_size):
                batch = data_rows[start:start + chunk_size]
                table_lines = []
                
                table_lines.append(f"### Sheet: {sheetname} (Rows {start + 1}-{start + len(batch)})")
                table_lines.append("| " + " | ".join(headers) + " |")
                table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
                
                for r in batch:
                    # Pad row to match header length
                    padded = r + [""] * (len(headers) - len(r))
                    table_lines.append("| " + " | ".join(padded[:len(headers)]) + " |")
                    
                units.append(ExtractedUnit(
                    index=unit_idx,
                    label=f"{sheetname} (R{start + 1}-{start + len(batch)})",
                    text="\n".join(table_lines)
                ))
                unit_idx += 1
                
        return units
    except Exception as e:
        logger.warning(f"openpyxl error on {filename}: {e}. Falling back to text extraction.")
        return extract_text_fallback(file_bytes, filename)


def extract_csv(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Extract CSV / TSV tabular data formatted as Markdown tables."""
    text_content = decode_bytes_to_string(file_bytes)
    delimiter = "\t" if filename.lower().endswith(".tsv") else ","
    
    try:
        reader = list(csv.reader(io.StringIO(text_content), delimiter=delimiter))
    except Exception:
        # Fallback to standard comma
        reader = list(csv.reader(io.StringIO(text_content)))
        
    if not reader:
        return []
        
    headers = [str(h).strip() for h in reader[0]]
    data_rows = reader[1:]
    
    units: List[ExtractedUnit] = []
    chunk_size = 30
    unit_idx = 1
    
    for start in range(0, max(len(data_rows), 1), chunk_size):
        batch = data_rows[start:start + chunk_size]
        table_lines = []
        
        table_lines.append(f"### Table Rows {start + 1} to {start + len(batch)}")
        table_lines.append("| " + " | ".join(headers) + " |")
        table_lines.append("| " + " | ".join(["---"] * len(headers)) + " |")
        
        for r in batch:
            padded = r + [""] * (len(headers) - len(r))
            table_lines.append("| " + " | ".join(padded[:len(headers)]) + " |")
            
        units.append(ExtractedUnit(
            index=unit_idx,
            label=f"Rows {start + 1}-{start + len(batch)}",
            text="\n".join(table_lines)
        ))
        unit_idx += 1
        
    return units


def extract_json(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Extract JSON / JSONL structured data into structured blocks."""
    text_content = decode_bytes_to_string(file_bytes)
    units: List[ExtractedUnit] = []
    
    # Check if JSONL (multiple JSON lines)
    lines = [l.strip() for l in text_content.splitlines() if l.strip()]
    if filename.lower().endswith(".jsonl") or (len(lines) > 1 and all(l.startswith("{") and l.endswith("}") for l in lines[:5])):
        chunk_size = 15
        unit_idx = 1
        for start in range(0, len(lines), chunk_size):
            batch = lines[start:start + chunk_size]
            formatted_records = []
            for i, line in enumerate(batch, start=start + 1):
                try:
                    obj = json.loads(line)
                    formatted_records.append(f"Record #{i}:\n" + json.dumps(obj, indent=2))
                except Exception:
                    formatted_records.append(f"Record #{i}:\n{line}")
                    
            units.append(ExtractedUnit(
                index=unit_idx,
                label=f"Records {start + 1}-{start + len(batch)}",
                text="\n\n".join(formatted_records)
            ))
            unit_idx += 1
        return units

    # Standard JSON object or array
    try:
        data = json.loads(text_content)
        if isinstance(data, list):
            chunk_size = 15
            unit_idx = 1
            for start in range(0, len(data), chunk_size):
                batch = data[start:start + chunk_size]
                text = f"### Records {start + 1} to {start + len(batch)}\n" + json.dumps(batch, indent=2)
                units.append(ExtractedUnit(
                    index=unit_idx,
                    label=f"Records {start + 1}-{start + len(batch)}",
                    text=text
                ))
                unit_idx += 1
            return units
        else:
            # Single JSON object: format with clean indentation
            formatted = json.dumps(data, indent=2)
            # Split into ~40 line sections if very large
            lines = formatted.splitlines()
            if len(lines) <= 50:
                return [ExtractedUnit(index=1, label="Document", text=formatted)]
                
            chunk_size = 40
            unit_idx = 1
            for start in range(0, len(lines), chunk_size):
                batch = lines[start:start + chunk_size]
                units.append(ExtractedUnit(
                    index=unit_idx,
                    label=f"Lines {start + 1}-{start + len(batch)}",
                    text="\n".join(batch)
                ))
                unit_idx += 1
            return units
    except Exception:
        return extract_text_fallback(file_bytes, filename)


def extract_code(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Extract source code sliced into line-number blocks (Lines 1-50, Lines 51-100)."""
    text_content = decode_bytes_to_string(file_bytes)
    lines = text_content.splitlines()
    if not lines:
        return []
        
    chunk_size = 45
    units: List[ExtractedUnit] = []
    unit_idx = 1
    
    ext = Path(filename).suffix.lstrip(".")
    
    for start in range(0, len(lines), chunk_size):
        batch = lines[start:start + chunk_size]
        numbered_lines = [f"{start + i + 1:4d} | {line}" for i, line in enumerate(batch)]
        
        block_text = f"```{ext}\n" + "\n".join(numbered_lines) + "\n```"
        units.append(ExtractedUnit(
            index=unit_idx,
            label=f"Lines {start + 1}-{start + len(batch)}",
            text=block_text
        ))
        unit_idx += 1
        
    return units


def extract_plain_text(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Extract plain text / markdown / rst / logs split into logical pages."""
    text_content = decode_bytes_to_string(file_bytes)
    paragraphs = text_content.split("\n\n")
    
    units: List[ExtractedUnit] = []
    current_page = []
    page_idx = 1
    word_count = 0
    
    for p in paragraphs:
        cleaned = p.strip()
        if not cleaned:
            continue
            
        current_page.append(cleaned)
        word_count += len(cleaned.split())
        
        if word_count >= 350:
            units.append(ExtractedUnit(
                index=page_idx,
                label=f"Page {page_idx}",
                text="\n\n".join(current_page)
            ))
            page_idx += 1
            current_page = []
            word_count = 0
            
    if current_page:
        units.append(ExtractedUnit(
            index=page_idx,
            label=f"Page {page_idx}",
            text="\n\n".join(current_page)
        ))
        
    if not units and text_content.strip():
        units.append(ExtractedUnit(index=1, label="Page 1", text=text_content.strip()))
        
    return units


def decode_bytes_to_string(file_bytes: bytes) -> str:
    """Losslessly decodes bytes using multiple encoding fallbacks."""
    for enc in ["utf-8", "utf-8-sig", "latin-1", "cp1252", "iso-8859-1"]:
        try:
            return file_bytes.decode(enc)
        except UnicodeDecodeError:
            continue
    return file_bytes.decode("utf-8", errors="replace")


def extract_text_fallback(file_bytes: bytes, filename: str) -> List[ExtractedUnit]:
    """Ultimate safe fallback that always succeeds."""
    text = decode_bytes_to_string(file_bytes)
    if not text.strip():
        return []
    return extract_plain_text(file_bytes, filename)


def extract_universal(file_bytes: bytes, filename: str) -> List[Tuple[int, str, str]]:
    """
    Main universal extraction dispatcher.
    Returns a list of (unit_index, unit_label, unit_text).
    """
    ext = Path(filename).suffix.lower()
    
    if ext == ".pdf":
        units = extract_pdf(file_bytes, filename)
    elif ext in [".docx", ".doc", ".odt", ".rtf"]:
        units = extract_docx(file_bytes, filename)
    elif ext in [".pptx", ".ppt"]:
        units = extract_pptx(file_bytes, filename)
    elif ext in [".xlsx", ".xls"]:
        units = extract_xlsx(file_bytes, filename)
    elif ext in [".csv", ".tsv"]:
        units = extract_csv(file_bytes, filename)
    elif ext in [".json", ".jsonl"]:
        units = extract_json(file_bytes, filename)
    elif ext in [".yaml", ".yml", ".xml", ".html", ".htm"]:
        units = extract_plain_text(file_bytes, filename)
    elif ext in [
        ".py", ".js", ".ts", ".tsx", ".jsx", ".java", ".cpp", ".c", ".h", ".hpp",
        ".go", ".rs", ".sql", ".sh", ".bash", ".zsh", ".css", ".scss", ".php", ".rb", ".swift", ".kt"
    ]:
        units = extract_code(file_bytes, filename)
    else:
        # Default text extractor
        units = extract_plain_text(file_bytes, filename)
        
    return [u.to_tuple() for u in units]
